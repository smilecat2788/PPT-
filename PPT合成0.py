from pptx import Presentation
from pptx.util import Inches
from pptx.util import Cm
import os
from PIL import Image


def rename():
    path= os.getcwd()
    filelist = os.listdir(path)  # 该文件夹下所有的文件（包括文件夹）
    for files in filelist:  # 遍历所有文件
        Olddir = os.path.join(path, files)  # 原来的文件路径
        if os.path.isdir(Olddir):  # 如果是文件夹则跳过
            continue
        filename = os.path.splitext(files)[0]  # 文件名
        filetype = os.path.splitext(files)[1]  # 文件扩展名

        if filename.startswith("PPT"):
            continue
        if filetype.endswith('png') or filetype.endswith('PNG'):
            newlen = len(filename)-3
            newfilename = filename[3:3+newlen]

            Newdir = os.path.join(path, newfilename+filetype)  # 新的文件路径
            os.rename(Olddir, Newdir)  # 重命名


i = 0
prs = Presentation()
prs.slide_height = Cm(19.05)
prs.slide_width = Cm(33.87)
title_slide_layout = prs.slide_layouts[6]  # 设置新建幻灯片样式为空白

files = os.listdir()  # 获取当前目录所有文件
#files.sort(key=lambda x: int(x.split('.')[0]))  # 对结果进行排序
for file in files:
    if os.path.splitext(file)[-1] == ".jpg" or os.path.splitext(file)[-1] == ".png" or os.path.splitext(file)[-1] == ".JPG" or os.path.splitext(file)[-1] == ".PNG":  # 判断是否是图片
        i = i+1
        print('正在生成第' + str(i) + '张PPT')
        slide = prs.slides.add_slide(title_slide_layout)  # 新建一页ppt
        pic = slide.shapes.add_picture(file, Inches(0), Inches(0))
print('生成完毕，共'+str(i)+'张PPT，正在保存')
prs.save('输出0.pptx')  # 创建PPT文件
print('保存完成')
