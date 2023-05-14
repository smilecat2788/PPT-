from pptx import Presentation
from pptx.util import Inches
from pptx.util import Cm
import os
from natsort import natsorted
i = 0
prs = Presentation()
prs.slide_height = Cm(19.05)
prs.slide_width = Cm(33.87)
title_slide_layout = prs.slide_layouts[6]  # 设置新建幻灯片样式为空白

files = os.listdir()  # 获取当前目录所有文件
files=natsorted(files)  # 对结果进行排序
#print(files)
for file in files:
    if os.path.splitext(file)[-1] == ".jpg" or os.path.splitext(file)[-1] == ".png" or os.path.splitext(file)[-1] == ".JPG" or os.path.splitext(file)[-1] == ".PNG":  # 判断是否是图片
        i = i+1
        print('正在生成第' + str(i) + '张PPT')
        slide = prs.slides.add_slide(title_slide_layout)  # 新建一页ppt
        pic = slide.shapes.add_picture(file, Inches(0), Inches(0))
print('生成完毕，共'+str(i)+'张PPT，正在保存')
prs.save('输出.pptx')  # 创建PPT文件
print('保存完成')
